"""Export QRHub Keynote slides to PowerPoint (.pptx)."""
import asyncio
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

BASE_URL = "https://qr-presentation-live.preview.emergentagent.com"
TOTAL_SLIDES = 30
OUT_DIR = Path("/app/scripts/slides_png")
OUT_PPTX = Path("/app/QRHub_Presentation.pptx")
WIDTH, HEIGHT = 1920, 1080


async def capture():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(viewport={"width": WIDTH, "height": HEIGHT})
        page = await ctx.new_page()
        for i in range(1, TOTAL_SLIDES + 1):
            url = f"{BASE_URL}/?slide={i}"
            print(f"[{i}/{TOTAL_SLIDES}] {url}")
            await page.goto(url, wait_until="networkidle", timeout=45000)
            await page.wait_for_timeout(2200)
            # Hide UI chrome (nav arrows, fullscreen btn, footer hints) before capture
            await page.add_style_tag(content="""
                button[aria-label],
                .slide-grain {opacity: 1}
                /* hide nav arrows and helper bottom bar if present */
                [class*='ChevronLeft'], [class*='ChevronRight'] {}
            """)
            out = OUT_DIR / f"slide_{i:02d}.png"
            await page.screenshot(path=str(out), full_page=False, type="png")
        await browser.close()


def build_pptx():
    prs = Presentation()
    # 16:9 slide size = 13.333" x 7.5" => EMU
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for i in range(1, TOTAL_SLIDES + 1):
        img = OUT_DIR / f"slide_{i:02d}.png"
        if not img.exists():
            print(f"Missing {img}, skipping")
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX}")


async def main():
    await capture()
    build_pptx()


if __name__ == "__main__":
    asyncio.run(main())
